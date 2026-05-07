"""PDFApps – TabConverter: convert PDF to images, DOCX, TXT, PPTX, XLSX, HTML, EPUB."""

import os
import re

_CTRL_RE = re.compile(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]')


def _clean(text: str) -> str:
    """Strip control characters that break XML-based formats."""
    return _CTRL_RE.sub('', text)

from PySide6.QtCore import Qt
from PySide6.QtWidgets import (
    QGroupBox, QFormLayout, QComboBox, QLabel, QFileDialog,
    QMessageBox, QProgressDialog,
)

from app.base import BasePage
from app.i18n import t
from app.utils import section, info_lbl, pick_folder
from app.constants import DESKTOP
from app.widgets import DropFileEdit


class TabConverter(BasePage):
    _DPI_VALUES = [72, 150, 300]

    def __init__(self, status_fn):
        super().__init__("fa5s.exchange-alt", t("tool.convert.name"),
                         t("tool.convert.desc"),
                         t("tool.convert.btn"), status_fn)
        f = self._form

        # -- Source file --
        sec_src = section(t("tool.convert.source"))
        f.addWidget(sec_src)
        self.drop_in = DropFileEdit()
        try: self.drop_in.btn.clicked.disconnect()
        except RuntimeError: pass
        self.drop_in.btn.clicked.connect(self._pick_input)
        self.drop_in.path_changed.connect(self._load_input)
        self.lbl_info = info_lbl()
        f.addWidget(self.drop_in); f.addWidget(self.lbl_info)

        # -- Output format --
        grp_fmt = QGroupBox(t("tool.convert.format_section"))
        gf = QFormLayout(grp_fmt)
        gf.setLabelAlignment(Qt.AlignmentFlag.AlignRight)
        self.cmb_format = QComboBox()
        self.cmb_format.addItems([
            t("tool.convert.png"), t("tool.convert.jpg"),
            t("tool.convert.docx"), t("tool.convert.txt"),
            t("tool.convert.pptx"), t("tool.convert.xlsx"),
            t("tool.convert.html"), t("tool.convert.epub"),
        ])
        self.cmb_format.currentIndexChanged.connect(self._on_format_changed)
        gf.addRow(t("tool.convert.format_label"), self.cmb_format)
        f.addWidget(grp_fmt)

        # -- Image options (visible for PNG/JPG) --
        self._grp_dpi = QGroupBox(t("tool.convert.img_options"))
        gd = QFormLayout(self._grp_dpi)
        gd.setLabelAlignment(Qt.AlignmentFlag.AlignRight)
        self.cmb_dpi = QComboBox()
        self.cmb_dpi.addItems([
            t("tool.convert.dpi_72"), t("tool.convert.dpi_150"), t("tool.convert.dpi_300"),
        ])
        self.cmb_dpi.setCurrentIndex(1)
        gd.addRow(t("tool.convert.dpi_label"), self.cmb_dpi)
        f.addWidget(self._grp_dpi)

        # -- Output folder (images) --
        sec_out_folder = section(t("tool.convert.output_folder"))
        f.addWidget(sec_out_folder)
        self._sec_out_folder = sec_out_folder
        self._drop_folder = DropFileEdit(placeholder=t("tool.convert.folder_hint"))
        try: self._drop_folder.btn.clicked.disconnect()
        except RuntimeError: pass
        self._drop_folder.btn.clicked.connect(self._pick_folder)
        f.addWidget(self._drop_folder)

        # -- Output file (DOCX / TXT) --
        self._section_file = section(t("tool.convert.output_file"))
        self._section_file.setVisible(False)
        f.addWidget(self._section_file)
        self._drop_file = DropFileEdit(save=True, default_name="converted.docx")
        self._drop_file.setVisible(False)
        f.addWidget(self._drop_file)

        self.lbl_result = QLabel("")
        self.lbl_result.setStyleSheet(
            "font-weight:600; font-size:11pt; color:#059669; "
            "background:transparent; padding:10px 4px;")
        f.addWidget(self.lbl_result)
        f.addStretch()
        self._compact_hidden = [sec_src, self.drop_in, self.lbl_info]
        # Hide all output sections — save dialog prompts automatically
        for w in (sec_out_folder, self._drop_folder, self._section_file, self._drop_file):
            w.setVisible(False)

    # ── UI callbacks ──────────────────────────────────────────────────────

    _EXT_MAP = {2: ".docx", 3: ".txt", 4: ".pptx", 5: ".xlsx", 6: ".html", 7: ".epub"}

    def _on_format_changed(self, index: int):
        is_image = index <= 1
        self._grp_dpi.setVisible(is_image)
        if not self._compact_active:
            self._drop_folder.setVisible(is_image)
            self._sec_out_folder.setVisible(is_image)
            self._section_file.setVisible(not is_image)
            self._drop_file.setVisible(not is_image)
        if not is_image:
            ext = self._EXT_MAP.get(index, ".pdf")
            inp = self.drop_in.path()
            if inp:
                base = os.path.splitext(inp)[0]
                self._drop_file.set_path(base + ext)

    def _pick_input(self):
        p, _ = QFileDialog.getOpenFileName(self, t("btn.open_pdf"), DESKTOP, t("file_filter.pdf"))
        if p:
            self._load_input(p)

    def _load_input(self, p: str):
        self.drop_in.blockSignals(True)
        self.drop_in.set_path(p)
        self.drop_in.blockSignals(False)
        if not self._maybe_prompt_password(p):
            self.drop_in.blockSignals(True); self.drop_in.set_path("")
            self.drop_in.blockSignals(False); return
        size = os.path.getsize(p)
        try:
            r = self._open_reader(p)
            self.lbl_info.setText(t("tool.compress.pages_info", n=len(r.pages), size=f"{size/1024:.1f}"))
        except Exception as e:
            self.lbl_info.setText(t("tool.split.error_info", e=e))
        # auto-set output paths
        base = os.path.splitext(p)[0]
        if not self._drop_folder.path():
            self._drop_folder.blockSignals(True)
            self._drop_folder.set_path(os.path.dirname(p))
            self._drop_folder.blockSignals(False)
        fmt = self.cmb_format.currentIndex()
        if fmt >= 2:
            ext = self._EXT_MAP.get(fmt, ".pdf")
            self._drop_file.set_path(base + ext)

    def _pick_folder(self):
        d = pick_folder(self)
        if d:
            self._drop_folder.blockSignals(True)
            self._drop_folder.set_path(d)
            self._drop_folder.blockSignals(False)

    def auto_load(self, path: str):
        if path and not self.drop_in.path():
            self._load_input(path)

    # ── conversion logic ──────────────────────────────────────────────────

    def _run(self):
        pdf_path = self.drop_in.path()
        if not pdf_path or not os.path.isfile(pdf_path):
            QMessageBox.warning(self, t("msg.warning"), t("msg.select_valid_pdf"))
            return

        fmt = self.cmb_format.currentIndex()
        self.lbl_result.setText("")

        converters = {
            0: lambda p: self._convert_images(p, 0),
            1: lambda p: self._convert_images(p, 1),
            2: self._convert_docx,
            3: self._convert_txt,
            4: self._convert_pptx,
            5: self._convert_xlsx,
            6: self._convert_html,
            7: self._convert_epub,
        }
        converters[fmt](pdf_path)

    def _make_progress(self, total: int, label: str) -> QProgressDialog:
        progress = QProgressDialog(label, t("progress.cancel"), 0, total, self)
        progress.setWindowTitle(t("progress.compress.title"))
        progress.setWindowModality(Qt.WindowModality.WindowModal)
        progress.setMinimumDuration(0)
        progress.setValue(0)
        return progress

    def _convert_images(self, pdf_path: str, fmt: int):
        out_dir = self._resolve_output_dir(self._drop_folder, pdf_path)
        if not out_dir:
            return
        os.makedirs(out_dir, exist_ok=True)
        ext = "png" if fmt == 0 else "jpg"
        dpi = self._DPI_VALUES[self.cmb_dpi.currentIndex()]
        self._status(f"→ {ext.upper()} @ {dpi} DPI…")

        try:
            with self._open_fitz(pdf_path) as _probe:
                total = _probe.page_count
        except Exception as e:
            QMessageBox.critical(self, t("msg.error"), str(e)); return

        pwd = self._pdf_password

        def do_work(worker):
            import fitz
            doc = fitz.open(pdf_path)
            if doc.needs_pass and pwd:
                doc.authenticate(pwd)
            try:
                matrix = fitz.Matrix(dpi / 72, dpi / 72)
                for i, page in enumerate(doc):
                    if worker.is_cancelled():
                        return None
                    worker.progress.emit(i, f"{i + 1}/{total}…")
                    pix = page.get_pixmap(matrix=matrix)
                    if pix.alpha:
                        pix = fitz.Pixmap(pix, 0)
                    if pix.n == 4:
                        pix = fitz.Pixmap(fitz.csRGB, pix)
                    out_file = os.path.join(out_dir, f"page_{i + 1:03d}.{ext}")
                    if ext == "png":
                        pix.save(out_file)
                    else:
                        try:
                            from PIL import Image
                            mode = "L" if pix.n == 1 else "RGB"
                            img = Image.frombytes(mode, (pix.width, pix.height), pix.samples)
                            img.save(out_file, "JPEG", quality=95)
                        except ImportError:
                            pix.save(out_file)
            finally:
                doc.close()
            return total

        def on_done(result):
            self.lbl_result.setText(f"  {result} → {out_dir}")
            self._status(f"✔  {result} images")
            QMessageBox.information(self, t("msg.done"),
                                    t("tool.convert.done_images",
                                      n=result, folder=out_dir))

        self._run_background(do_work, total, t("tool.convert.converting"),
                             on_done=on_done)

    def _convert_docx(self, pdf_path: str):
        out_path = self._resolve_output_file(self._drop_file, pdf_path,
                                             filter_key="file_filter.docx")
        if not out_path:
            return
        # Pre-flight on main thread: dep checks + page count + capture pwd.
        try:
            import fitz  # noqa: F401
        except ImportError:
            QMessageBox.critical(self, t("msg.missing_dep"), t("tool.ocr.dep_pymupdf"))
            return
        try:
            from docx import Document  # noqa: F401
        except ImportError:
            QMessageBox.critical(self, t("msg.missing_dep"), t("tool.convert.dep_docx"))
            return
        try:
            with self._open_fitz(pdf_path) as _probe:
                total = _probe.page_count
        except Exception as e:
            QMessageBox.critical(self, t("msg.error"), str(e))
            return
        if total == 0:
            QMessageBox.warning(self, t("msg.warning"), t("msg.select_valid_pdf"))
            return
        pwd = self._pdf_password

        def do_work(worker):
            import fitz
            from docx import Document
            from docx.shared import Pt, RGBColor, Inches
            import io, re as _re
            doc = fitz.open(pdf_path)
            if doc.needs_pass and pwd:
                doc.authenticate(pwd)
            try:
                docx_doc = Document()
                for i, page in enumerate(doc):
                    if worker.is_cancelled():
                        return None
                    blocks = page.get_text("dict")["blocks"]
                    for block in blocks:
                        btype = block.get("type", 0)
                        # Image block — extract and embed
                        if btype == 1:
                            img_data = block.get("image")
                            if img_data:
                                try:
                                    para = docx_doc.add_paragraph()
                                    run = para.add_run()
                                    run.add_picture(io.BytesIO(img_data), width=Inches(5.0))
                                except Exception:
                                    pass
                            continue
                        # Text block
                        lines = block.get("lines", [])
                        if not lines:
                            continue
                        all_spans = []
                        for line in lines:
                            all_spans.extend(line.get("spans", []))
                        if not all_spans:
                            continue
                        block_text = " ".join(
                            _clean(s.get("text", "")) for s in all_spans
                        ).strip()
                        if not block_text:
                            continue
                        # Skip standalone page numbers
                        if _re.match(r"^\s*(?:page\s+)?\d{1,4}(?:\s+of\s+\d{1,4})?\s*$",
                                     block_text, _re.IGNORECASE):
                            continue
                        # Skip TOC dot-leader lines
                        if _re.search(r'\.[\s.]*\.[\s.]*\.[\s.]*\.', block_text):
                            continue
                        # Detect heading level by font size
                        max_size = max(s.get("size", 12) for s in all_spans)
                        any_bold = any(s.get("flags", 0) & 16 for s in all_spans)
                        if max_size >= 20:
                            para = docx_doc.add_heading(level=1)
                        elif max_size >= 16:
                            para = docx_doc.add_heading(level=2)
                        elif max_size >= 13 and any_bold:
                            para = docx_doc.add_heading(level=3)
                        elif any_bold and max_size >= 11:
                            para = docx_doc.add_heading(level=4)
                        else:
                            para = docx_doc.add_paragraph()
                        for li, line in enumerate(lines):
                            spans = line.get("spans", [])
                            for span in spans:
                                text = _clean(span.get("text", ""))
                                if not text:
                                    continue
                                run = para.add_run(text)
                                run.font.size = Pt(span.get("size", 12))
                                sf = span.get("flags", 0)
                                run.font.bold = bool(sf & 16)
                                run.font.italic = bool(sf & 2)
                                color = span.get("color", 0)
                                if color and color != 0:
                                    r_val = (color >> 16) & 0xFF
                                    g_val = (color >> 8) & 0xFF
                                    b_val = color & 0xFF
                                    run.font.color.rgb = RGBColor(r_val, g_val, b_val)
                            if li < len(lines) - 1:
                                para.add_run(" ")
                    worker.progress.emit(i, f"{i + 1}/{total}…")
                if worker.is_cancelled():
                    return None
                docx_doc.save(out_path)
            finally:
                doc.close()
            return total

        def on_done(result):
            self.lbl_result.setText(f"  → {os.path.basename(out_path)}")
            self._status(f"✔  DOCX → {out_path}")
            QMessageBox.information(self, t("msg.done"),
                                    t("tool.convert.done_docx", path=out_path))

        self._run_background(do_work, total, t("tool.convert.converting"),
                             on_done=on_done)

    def _convert_txt(self, pdf_path: str):
        out_path = self._resolve_output_file(self._drop_file, pdf_path,
                                             filter_key="file_filter.txt")
        if not out_path:
            return
        try:
            with self._open_fitz(pdf_path) as _probe:
                total = _probe.page_count
        except Exception as e:
            QMessageBox.critical(self, t("msg.error"), str(e))
            return
        if total == 0:
            QMessageBox.warning(self, t("msg.warning"), t("msg.select_valid_pdf"))
            return
        pwd = self._pdf_password

        def do_work(worker):
            import fitz
            doc = fitz.open(pdf_path)
            if doc.needs_pass and pwd:
                doc.authenticate(pwd)
            try:
                with open(out_path, 'w', encoding='utf-8') as f:
                    for i, page in enumerate(doc):
                        if worker.is_cancelled():
                            return None
                        if i > 0:
                            f.write(f'\n\n--- Page {i + 1} ---\n\n')
                        f.write(page.get_text())
                        worker.progress.emit(i, f"{i + 1}/{total}…")
            finally:
                doc.close()
            return total

        def on_done(result):
            self.lbl_result.setText(f"  → {os.path.basename(out_path)}")
            self._status(f"✔  TXT → {out_path}")
            QMessageBox.information(self, t("msg.done"),
                                    t("tool.convert.done_txt", path=out_path))

        self._run_background(do_work, total, t("tool.convert.converting"),
                             on_done=on_done)

    # ── PDF → PPTX ──────────────────────────────────────────────────────

    def _convert_pptx(self, pdf_path: str):
        out_path = self._resolve_output_file(self._drop_file, pdf_path,
                                             filter_key="file_filter.pptx")
        if not out_path:
            return
        try:
            from pptx import Presentation  # noqa: F401
        except ImportError:
            QMessageBox.critical(self, t("msg.missing_dep"), t("tool.convert.dep_pptx"))
            return
        try:
            with self._open_fitz(pdf_path) as _probe:
                total = _probe.page_count
                first = _probe[0].rect if total else None
        except Exception as e:
            QMessageBox.critical(self, t("msg.error"), str(e))
            return
        if total == 0:
            QMessageBox.warning(self, t("msg.warning"), t("msg.select_valid_pdf"))
            return
        slide_w_pt = first.width
        slide_h_pt = first.height
        pwd = self._pdf_password

        def do_work(worker):
            import fitz, io
            from pptx import Presentation
            from pptx.util import Emu, Pt
            from pptx.dml.color import RGBColor
            doc = fitz.open(pdf_path)
            if doc.needs_pass and pwd:
                doc.authenticate(pwd)
            try:
                prs = Presentation()
                prs.slide_width = Emu(int(slide_w_pt * 12700))
                prs.slide_height = Emu(int(slide_h_pt * 12700))
                blank = prs.slide_layouts[6]
                for i, page in enumerate(doc):
                    if worker.is_cancelled():
                        return None
                    slide = prs.slides.add_slide(blank)

                    # Extract text blocks and image blocks separately so
                    # the result is editable in PowerPoint instead of
                    # one rasterised image per slide. Each text block
                    # becomes a textbox at its bbox; each line in the
                    # block becomes a paragraph; each span becomes a
                    # run with its own size/bold/italic/color.
                    blocks = page.get_text("dict").get("blocks", [])
                    for block in blocks:
                        bbox = block.get("bbox")
                        if not bbox:
                            continue
                        x0, y0, x1, y1 = bbox
                        # Clamp to slide bounds; fitz can emit blocks
                        # whose bbox slightly overflows the page rect
                        # for italic/decorated glyphs.
                        x0 = max(0, x0); y0 = max(0, y0)
                        x1 = min(slide_w_pt, x1); y1 = min(slide_h_pt, y1)
                        w = x1 - x0; h = y1 - y0
                        if w <= 0 or h <= 0:
                            continue

                        if block.get("type") == 1:  # image block
                            img_data = block.get("image")
                            if not img_data:
                                continue
                            try:
                                slide.shapes.add_picture(
                                    io.BytesIO(img_data),
                                    Emu(int(x0 * 12700)),
                                    Emu(int(y0 * 12700)),
                                    Emu(int(w * 12700)),
                                    Emu(int(h * 12700)))
                            except Exception:
                                # Image format not supported by python-pptx
                                # (rare formats like JBIG2). Skip the
                                # block — the rest of the slide is still
                                # produced.
                                pass
                            continue

                        lines = block.get("lines", [])
                        if not lines:
                            continue
                        try:
                            tb = slide.shapes.add_textbox(
                                Emu(int(x0 * 12700)),
                                Emu(int(y0 * 12700)),
                                Emu(int(w * 12700)),
                                Emu(int(h * 12700)))
                        except Exception:
                            continue
                        tf = tb.text_frame
                        tf.word_wrap = True
                        # Zero internal padding so the textbox bbox
                        # matches the PDF span bbox more faithfully.
                        # Older python-pptx versions reject Emu(0) for
                        # margins; tolerate that and leave the default.
                        for attr in ("margin_left", "margin_right",
                                     "margin_top", "margin_bottom"):
                            try: setattr(tf, attr, 0)
                            except Exception: pass  # noqa: S110

                        first_line = True
                        for line in lines:
                            spans = line.get("spans", [])
                            if not spans:
                                continue
                            para = tf.paragraphs[0] if first_line else tf.add_paragraph()
                            first_line = False
                            first_run = True
                            for span in spans:
                                text = span.get("text", "")
                                if not text:
                                    continue
                                # Reuse the auto-created empty run for
                                # the first span; add new runs after.
                                if first_run and len(para.runs) > 0:
                                    run = para.runs[0]
                                else:
                                    run = para.add_run()
                                first_run = False
                                run.text = text
                                size = span.get("size", 12)
                                # Pt() rejects negative or non-numeric
                                # sizes; if the PDF span had a junk
                                # value, fall back to the theme default.
                                try: run.font.size = Pt(size)
                                except Exception: pass  # noqa: S110
                                flags = span.get("flags", 0)
                                run.font.bold = bool(flags & 16)
                                run.font.italic = bool(flags & 2)
                                color = span.get("color", 0)
                                if color:
                                    try:
                                        run.font.color.rgb = RGBColor(
                                            (color >> 16) & 0xFF,
                                            (color >> 8) & 0xFF,
                                            color & 0xFF)
                                    except Exception:
                                        # Some run types reject explicit
                                        # color (e.g. inside placeholder
                                        # layouts); leave the default.
                                        pass
                    worker.progress.emit(i, f"{i + 1}/{total}…")
                if worker.is_cancelled():
                    return None
                prs.save(out_path)
            finally:
                doc.close()
            return total

        def on_done(result):
            self.lbl_result.setText(f"  → {os.path.basename(out_path)}")
            self._status(f"✔  PPTX → {out_path}")
            QMessageBox.information(self, t("msg.done"),
                                    t("tool.convert.done_pptx", path=out_path))

        self._run_background(do_work, total, t("tool.convert.converting"),
                             on_done=on_done)

    # ── PDF → XLSX ──────────────────────────────────────────────────────

    def _convert_xlsx(self, pdf_path: str):
        out_path = self._resolve_output_file(self._drop_file, pdf_path,
                                             filter_key="file_filter.xlsx")
        if not out_path:
            return
        try:
            from openpyxl import Workbook  # noqa: F401
        except ImportError:
            QMessageBox.critical(self, t("msg.missing_dep"), t("tool.convert.dep_xlsx"))
            return
        try:
            with self._open_fitz(pdf_path) as _probe:
                total = _probe.page_count
        except Exception as e:
            QMessageBox.critical(self, t("msg.error"), str(e))
            return
        if total == 0:
            QMessageBox.warning(self, t("msg.warning"), t("msg.select_valid_pdf"))
            return
        pwd = self._pdf_password

        def do_work(worker):
            import fitz
            from openpyxl import Workbook
            doc = fitz.open(pdf_path)
            if doc.needs_pass and pwd:
                doc.authenticate(pwd)
            try:
                wb = Workbook()
                wb.remove(wb.active)
                for i, page in enumerate(doc):
                    if worker.is_cancelled():
                        return None
                    ws = wb.create_sheet(title=f"Page {i + 1}")
                    blocks = page.get_text("blocks")
                    for row_idx, block in enumerate(blocks):
                        if block[6] != 0:  # skip image blocks
                            continue
                        text = _clean(block[4].strip())
                        if text:
                            cells = [c.strip() for c in text.replace("\t", "|").split("|") if c.strip()]
                            if not cells:
                                cells = [text]
                            for col_idx, cell in enumerate(cells):
                                ws.cell(row=row_idx + 1, column=col_idx + 1, value=cell)
                    worker.progress.emit(i, f"{i + 1}/{total}…")
                if worker.is_cancelled():
                    return None
                wb.save(out_path)
            finally:
                doc.close()
            return total

        def on_done(result):
            self.lbl_result.setText(f"  → {os.path.basename(out_path)}")
            self._status(f"✔  XLSX → {out_path}")
            QMessageBox.information(self, t("msg.done"),
                                    t("tool.convert.done_xlsx", path=out_path))

        self._run_background(do_work, total, t("tool.convert.converting"),
                             on_done=on_done)

    # ── PDF → HTML ──────────────────────────────────────────────────────

    def _convert_html(self, pdf_path: str):
        out_path = self._resolve_output_file(self._drop_file, pdf_path,
                                             filter_key="file_filter.html")
        if not out_path:
            return
        try:
            with self._open_fitz(pdf_path) as _probe:
                total = _probe.page_count
        except Exception as e:
            QMessageBox.critical(self, t("msg.error"), str(e))
            return
        if total == 0:
            QMessageBox.warning(self, t("msg.warning"), t("msg.select_valid_pdf"))
            return
        pwd = self._pdf_password

        def do_work(worker):
            import fitz
            doc = fitz.open(pdf_path)
            if doc.needs_pass and pwd:
                doc.authenticate(pwd)
            try:
                parts = [
                    "<!DOCTYPE html>",
                    '<html lang="en"><head><meta charset="UTF-8">',
                    f"<title>{os.path.basename(pdf_path)}</title>",
                    "<style>body{font-family:sans-serif;max-width:800px;margin:0 auto;padding:20px;}"
                    ".page{margin-bottom:40px;padding-bottom:20px;border-bottom:1px solid #ccc;}</style>",
                    "</head><body>",
                ]
                for i, page in enumerate(doc):
                    if worker.is_cancelled():
                        return None
                    parts.append('<div class="page">')
                    blocks = page.get_text("dict")["blocks"]
                    for block in blocks:
                        if block.get("type") != 0:
                            continue
                        for line in block.get("lines", []):
                            spans_html = ""
                            for span in line.get("spans", []):
                                text = span["text"]
                                if not text.strip():
                                    continue
                                flags = span.get("flags", 0)
                                bold = flags & 16
                                italic = flags & 2
                                tag_text = _clean(text).replace("&", "&amp;").replace("<", "&lt;")
                                if bold:
                                    tag_text = f"<strong>{tag_text}</strong>"
                                if italic:
                                    tag_text = f"<em>{tag_text}</em>"
                                spans_html += tag_text
                            if spans_html:
                                avg_size = max(s.get("size", 12) for s in line.get("spans", [{"size": 12}]))
                                if avg_size >= 18:
                                    parts.append(f"<h1>{spans_html}</h1>")
                                elif avg_size >= 15:
                                    parts.append(f"<h2>{spans_html}</h2>")
                                elif avg_size >= 13:
                                    parts.append(f"<h3>{spans_html}</h3>")
                                else:
                                    parts.append(f"<p>{spans_html}</p>")
                    parts.append("</div>")
                    worker.progress.emit(i, f"{i + 1}/{total}…")
                parts.append("</body></html>")
                if worker.is_cancelled():
                    return None
                with open(out_path, "w", encoding="utf-8") as f:
                    f.write("\n".join(parts))
            finally:
                doc.close()
            return total

        def on_done(result):
            self.lbl_result.setText(f"  → {os.path.basename(out_path)}")
            self._status(f"✔  HTML → {out_path}")
            QMessageBox.information(self, t("msg.done"),
                                    t("tool.convert.done_html", path=out_path))

        self._run_background(do_work, total, t("tool.convert.converting"),
                             on_done=on_done)

    # ── PDF → EPUB ──────────────────────────────────────────────────────

    def _convert_epub(self, pdf_path: str):
        out_path = self._resolve_output_file(self._drop_file, pdf_path,
                                             filter_key="file_filter.epub")
        if not out_path:
            return
        try:
            from ebooklib import epub  # noqa: F401
        except ImportError:
            QMessageBox.critical(self, t("msg.missing_dep"), t("tool.convert.dep_epub"))
            return
        try:
            with self._open_fitz(pdf_path) as _probe:
                total = _probe.page_count
        except Exception as e:
            QMessageBox.critical(self, t("msg.error"), str(e))
            return
        if total == 0:
            QMessageBox.warning(self, t("msg.warning"), t("msg.select_valid_pdf"))
            return
        pwd = self._pdf_password

        def do_work(worker):
            import fitz
            from ebooklib import epub
            doc = fitz.open(pdf_path)
            if doc.needs_pass and pwd:
                doc.authenticate(pwd)
            try:
                book = epub.EpubBook()
                book.set_identifier(f"pdfapps-{os.path.basename(pdf_path)}")
                book.set_title(os.path.splitext(os.path.basename(pdf_path))[0])
                book.set_language("en")
                chapters = []
                for i, page in enumerate(doc):
                    if worker.is_cancelled():
                        return None
                    ch = epub.EpubHtml(title=f"Page {i + 1}", file_name=f"page_{i+1}.xhtml")
                    text = page.get_text()
                    paragraphs = [f"<p>{_clean(p)}</p>" for p in text.split("\n") if p.strip()]
                    ch.content = f"<html><body><h2>Page {i + 1}</h2>{''.join(paragraphs)}</body></html>"
                    book.add_item(ch)
                    chapters.append(ch)
                    worker.progress.emit(i, f"{i + 1}/{total}…")
                if worker.is_cancelled():
                    return None
                book.add_item(epub.EpubNcx())
                book.add_item(epub.EpubNav())
                book.spine = ["nav"] + chapters
                book.toc = chapters
                epub.write_epub(out_path, book)
            finally:
                doc.close()
            return total

        def on_done(result):
            self.lbl_result.setText(f"  → {os.path.basename(out_path)}")
            self._status(f"✔  EPUB → {out_path}")
            QMessageBox.information(self, t("msg.done"),
                                    t("tool.convert.done_epub", path=out_path))

        self._run_background(do_work, total, t("tool.convert.converting"),
                             on_done=on_done)
